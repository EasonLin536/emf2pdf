from pptx import Presentation 
from PyPDF2 import PdfFileReader, PdfFileWriter
import os
import sys

exec_path = os.getcwd()
print(exec_path)

if not os.path.isdir('{}/tmp'.format(exec_path)):
    os.system('mkdir {}/tmp'.format(exec_path))

# Creating an pptx file
img_path = '{}/{}'.format(exec_path, sys.argv[1]) 
file_prefix = sys.argv[1][:len(sys.argv[1])-4]
ppt = Presentation() 
blank_slide_layout = ppt.slide_layouts[6] 
slide = ppt.slides.add_slide(blank_slide_layout) 
pic = slide.shapes.add_picture(img_path, 0, 0)
(pic_width, pic_height) = pic.image.size
ppt.save('{}/tmp/tmp.pptx'.format(exec_path))

# Convert to pdf
os.system('''
    cd tmp/
    libreoffice --headless --invisible --convert-to pdf tmp.pptx
    cd ..''')

# Crop pdf
in_f = open("{}/tmp/tmp.pdf".format(exec_path), "rb")
page = PdfFileReader(in_f).getPage(0)
page_width = page.cropBox.getUpperRight_x()
page_height = page.cropBox.getUpperRight_y()
page.cropBox.lowerLeft = (0, page_height - pic_height * 0.75)
page.cropBox.upperRight = (pic_width * 0.75, page_height)
out_file = PdfFileWriter()
out_file.addPage(page)

with open("{}/{}.pdf".format(exec_path, file_prefix), "wb") as out_f:
    out_file.write(out_f)

in_f.close()
out_f.close()